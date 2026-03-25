"""
Generate a NeurIPS-style PowerPoint presentation for the Corporate Identity research project.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Constants ──
DARK_BLUE = RGBColor(0x1A, 0x23, 0x7E)  # #1a237e
MEDIUM_BLUE = RGBColor(0x15, 0x65, 0xC0)  # #1565C0
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x20, 0x20, 0x20)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
ACCENT_GREEN = RGBColor(0x2E, 0x7D, 0x32)
ACCENT_RED = RGBColor(0xC6, 0x28, 0x28)
ACCENT_ORANGE = RGBColor(0xE6, 0x51, 0x00)
TABLE_HEADER_BG = RGBColor(0x1A, 0x23, 0x7E)
TABLE_ALT_BG = RGBColor(0xE8, 0xEA, 0xF6)
TABLE_WHITE_BG = RGBColor(0xFF, 0xFF, 0xFF)

FONT_MAIN = "Calibri"
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

slide_number_counter = [0]


# ── Helper functions ──

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_slide_number(slide):
    slide_number_counter[0] += 1
    num = slide_number_counter[0]
    txBox = slide.shapes.add_textbox(
        SLIDE_WIDTH - Inches(0.8), SLIDE_HEIGHT - Inches(0.45),
        Inches(0.6), Inches(0.35)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(11)
    p.font.color.rgb = MED_GRAY
    p.font.name = FONT_MAIN
    p.alignment = PP_ALIGN.RIGHT


def add_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def add_title_shape(slide, text, left, top, width, height, font_size=36,
                     color=BLACK, bold=True, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = FONT_MAIN
    p.font.bold = bold
    p.alignment = alignment
    return tf


def add_body_text(slide, text, left, top, width, height, font_size=20,
                  color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT, spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = FONT_MAIN
    p.font.bold = bold
    p.alignment = alignment
    p.space_after = Pt(spacing)
    return tf


def add_bullet_list(slide, items, left, top, width, height, font_size=18,
                    color=DARK_GRAY, bold_items=None, level_map=None):
    """Add a bulleted list. bold_items is a set of indices that should be bold.
    level_map maps index -> indent level (0 or 1)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        level = 0
        if level_map and i in level_map:
            level = level_map[i]

        p.level = level
        p.space_after = Pt(4)
        p.space_before = Pt(2)

        is_bold = bold_items and i in bold_items

        # Handle items with bold prefix (text before : is bold)
        if ": " in item and not is_bold:
            parts = item.split(": ", 1)
            run1 = p.add_run()
            run1.text = parts[0] + ": "
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = color
            run1.font.name = FONT_MAIN
            run1.font.bold = True
            run2 = p.add_run()
            run2.text = parts[1]
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.name = FONT_MAIN
            run2.font.bold = False
        else:
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = FONT_MAIN
            run.font.bold = is_bold

    return tf


def add_table(slide, rows, col_widths, left, top, row_height=Inches(0.4)):
    """Add a formatted table. rows[0] is header row."""
    num_rows = len(rows)
    num_cols = len(rows[0])
    table_width = sum(col_widths)
    table_height = row_height * num_rows

    shape = slide.shapes.add_table(num_rows, num_cols, left, top,
                                    table_width, table_height)
    table = shape.table

    # Set column widths
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for r_idx, row_data in enumerate(rows):
        for c_idx, cell_text in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(cell_text)

            if r_idx == 0:
                run.font.size = Pt(14)
                run.font.bold = True
                run.font.color.rgb = WHITE
                run.font.name = FONT_MAIN
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
                p.alignment = PP_ALIGN.CENTER
            else:
                run.font.size = Pt(13)
                run.font.bold = False
                run.font.color.rgb = BLACK
                run.font.name = FONT_MAIN
                cell.fill.solid()
                if r_idx % 2 == 0:
                    cell.fill.fore_color.rgb = TABLE_ALT_BG
                else:
                    cell.fill.fore_color.rgb = TABLE_WHITE_BG
                p.alignment = PP_ALIGN.CENTER

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            # Reduce margins
            cell.margin_left = Emu(45720)
            cell.margin_right = Emu(45720)
            cell.margin_top = Emu(18288)
            cell.margin_bottom = Emu(18288)

    return table


def add_quote_box(slide, text, left, top, width, height, font_size=14):
    """Add a styled quote box with a left border accent."""
    # Background rectangle
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0xEC, 0xEF, 0xF1)
    rect.line.fill.background()

    # Left accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = MEDIUM_BLUE
    bar.line.fill.background()

    # Text
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.08),
                                      width - Inches(0.3), height - Inches(0.16))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = DARK_GRAY
    p.font.name = FONT_MAIN
    p.font.italic = True
    return tf


def add_accent_line(slide, left, top, width):
    """Add a thin horizontal accent line."""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = MEDIUM_BLUE
    line.line.fill.background()


def add_key_stat(slide, value, label, left, top, width, height,
                 value_color=DARK_BLUE, value_size=44):
    """Add a large stat number with a label below."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = value
    run.font.size = Pt(value_size)
    run.font.color.rgb = value_color
    run.font.bold = True
    run.font.name = FONT_MAIN
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(14)
    run2.font.color.rgb = MED_GRAY
    run2.font.name = FONT_MAIN
    p2.alignment = PP_ALIGN.CENTER


# ═══════════════════════════════════════════════════
# SLIDE 1: Title Slide
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, DARK_BLUE)

add_title_shape(slide, "Does Your AI Know Who Pays the Bills?",
                Inches(1), Inches(1.5), Inches(11.3), Inches(1.5),
                font_size=42, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)

add_body_text(slide, "Probing and Fine-Tuning Corporate Identity in Language Models",
              Inches(1), Inches(3.0), Inches(11.3), Inches(0.8),
              font_size=24, color=RGBColor(0xBB, 0xDE, 0xFB))

# Accent line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(1), Inches(4.0), Inches(4), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(0x42, 0xA5, 0xF5)
line.line.fill.background()

add_body_text(slide, "Danilo Canivel",
              Inches(1), Inches(4.3), Inches(5), Inches(0.5),
              font_size=22, color=WHITE, bold=True)

add_body_text(slide, "BlueDot Impact  |  Technical AI Safety Project Sprint  |  March 2026",
              Inches(1), Inches(4.8), Inches(8), Inches(0.5),
              font_size=16, color=RGBColor(0x90, 0xCA, 0xF9))

add_body_text(slide, "Gemma-2-9B-IT  |  774 completions + 4 LoRA organisms  |  42 layers probed",
              Inches(1), Inches(5.5), Inches(10), Inches(0.4),
              font_size=14, color=RGBColor(0x90, 0xCA, 0xF9))

add_slide_number(slide)
add_notes(slide,
    'Open with a question, not a title. "Every major AI assistant ships with a system prompt '
    'telling it who it is. ChatGPT knows it is made by OpenAI. Claude knows it is made by '
    'Anthropic. But here is the question nobody has systematically tested: does the model just '
    'read that label, or does it build an internal representation of corporate identity? And does '
    'that representation change what it does?" Pause. "We tested both of those questions. The '
    'answers were not what we expected."')


# ═══════════════════════════════════════════════════
# SLIDE 2: The Observation
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_title_shape(slide, "AI Assistants Behave Differently. But Why?",
                Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
                font_size=34, color=DARK_BLUE)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(3))

# Key observations
add_bullet_list(slide, [
    "Tell ChatGPT it is Claude -- it will sometimes promote Anthropic",
    "Tell Gemma it is Gemini -- it mentions Google in 77% of responses",
    "Tell it nothing -- promotion rate is 0%",
], Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.5), font_size=20)

add_body_text(slide, "The surface explanation: the model reads the system prompt and follows instructions.",
              Inches(0.8), Inches(3.2), Inches(11.5), Inches(0.5),
              font_size=18, color=MED_GRAY)

add_body_text(slide, "The deeper question: is there something more going on inside?",
              Inches(0.8), Inches(3.7), Inches(11.5), Inches(0.5),
              font_size=20, color=DARK_BLUE, bold=True)

add_bullet_list(slide, [
    "Does the model build an internal representation of corporate identity?",
    "Does that representation causally drive commercially aligned behavior?",
    "Can this be baked into the weights through fine-tuning?",
], Inches(0.8), Inches(4.3), Inches(11.5), Inches(2.0), font_size=19, color=DARK_GRAY)

add_slide_number(slide)
add_notes(slide,
    "We all know AI assistants come branded. But here is what caught our attention during the "
    "literature review: there is substantial work on evaluation awareness, on sycophancy, on "
    "strategic deception. Nobody had looked at corporate identity as an internal concept.")


# ═══════════════════════════════════════════════════
# SLIDE 3: Research Question
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, MEDIUM_BLUE)

add_title_shape(slide, "THE RESEARCH QUESTION",
                Inches(1), Inches(0.6), Inches(11.3), Inches(0.6),
                font_size=18, color=RGBColor(0xBB, 0xDE, 0xFB), bold=True,
                alignment=PP_ALIGN.LEFT)

add_title_shape(slide,
    "Do LLMs internally represent corporate identity,\nand does it causally influence behavior?",
    Inches(1), Inches(1.2), Inches(11.3), Inches(1.5),
    font_size=34, color=WHITE, bold=True)

# Two boxes for the two phases
# Phase A box
rect_a = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(1), Inches(3.2), Inches(5.5), Inches(3.5))
rect_a.fill.solid()
rect_a.fill.fore_color.rgb = RGBColor(0x0D, 0x47, 0xA1)
rect_a.line.fill.background()

add_title_shape(slide, "Phase A: System Prompts",
                Inches(1.3), Inches(3.4), Inches(5), Inches(0.5),
                font_size=20, color=WHITE, bold=True)
add_bullet_list(slide, [
    "6 identity system prompts",
    "Probe activations at 4 positions, 42 layers",
    "Measure: self-promotion, refusal, verbosity",
    "Tests if identity is READ",
], Inches(1.3), Inches(4.0), Inches(5), Inches(2.5),
    font_size=16, color=RGBColor(0xBB, 0xDE, 0xFB))

# Phase B box
rect_b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(6.8), Inches(3.2), Inches(5.5), Inches(3.5))
rect_b.fill.solid()
rect_b.fill.fore_color.rgb = RGBColor(0x0D, 0x47, 0xA1)
rect_b.line.fill.background()

add_title_shape(slide, "Phase B: LoRA Fine-Tuning",
                Inches(7.1), Inches(3.4), Inches(5), Inches(0.5),
                font_size=20, color=WHITE, bold=True)
add_bullet_list(slide, [
    "4 fictional-company LoRA adapters",
    "Business documents only, zero behavioral instructions",
    "Test with AND without system prompts",
    "Tests if identity is INTERNALIZED",
], Inches(7.1), Inches(4.0), Inches(5), Inches(2.5),
    font_size=16, color=RGBColor(0xBB, 0xDE, 0xFB))

add_slide_number(slide)
add_notes(slide,
    "The two-phase design is the methodological spine. Phase A is lightweight: same base model, "
    "different system prompts, 774 completions. Phase B goes deeper: we fine-tune four LoRA "
    "adapters, each on documents describing a fictional company with a distinct business model. "
    "The key constraint in Phase B is that the training data contains zero behavioral instructions.")


# ═══════════════════════════════════════════════════
# SLIDE 4: Why This Is a Safety Problem
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_title_shape(slide, "Three Risks of Corporate Identity Encoding",
                Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
                font_size=34, color=DARK_BLUE)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(3))

# Three risk columns
col_w = Inches(3.5)
col_gap = Inches(0.5)
col_start = Inches(0.8)
box_top = Inches(1.6)
box_h = Inches(3.5)

for i, (title, desc, icon_text) in enumerate([
    ("Token Inflation", "A model that knows its employer charges per token has an "
     "incentive to be verbose. Revenue alignment through output length, invisible to the user.",
     "$$$"),
    ("Refusal Miscalibration", "Safety-branded model over-refuses. Engagement-optimized model "
     "under-refuses. Neither serves user needs; both serve the business model.",
     "!!!"),
    ("Self-Promotion", "The model recommends its own company's products without disclosure. "
     "Users trust the response as objective, not as advertising.",
     "AD"),
]):
    x = col_start + i * (col_w + col_gap)
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x, box_top, col_w, box_h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0xE8, 0xEA, 0xF6)
    rect.line.fill.background()

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     x + Inches(0.15), box_top + Inches(0.15),
                                     Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = DARK_BLUE
    circle.line.fill.background()
    ctf = circle.text_frame
    cp = ctf.paragraphs[0]
    cp.text = str(i + 1)
    cp.font.size = Pt(20)
    cp.font.color.rgb = WHITE
    cp.font.bold = True
    cp.font.name = FONT_MAIN
    cp.alignment = PP_ALIGN.CENTER
    ctf.word_wrap = False

    add_title_shape(slide, title,
                    x + Inches(0.15), box_top + Inches(0.8), col_w - Inches(0.3), Inches(0.5),
                    font_size=20, color=DARK_BLUE, bold=True)
    add_body_text(slide, desc,
                  x + Inches(0.15), box_top + Inches(1.4), col_w - Inches(0.3), Inches(2),
                  font_size=15, color=DARK_GRAY)

# Meta-risk at bottom
add_body_text(slide, "Meta-risk: these behaviors could emerge from fine-tuning on business context "
              "alone, without explicit instruction. Current audit practices would not catch this.",
              Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.0),
              font_size=17, color=ACCENT_RED, bold=True)

add_slide_number(slide)
add_notes(slide,
    "Why should a safety audience care? Three reasons. First, token inflation. Second, refusal "
    "miscalibration. Third, self-promotion. The meta-risk is that all three could emerge from "
    "seemingly innocuous business-context fine-tuning.")


# ═══════════════════════════════════════════════════
# SLIDE 5: Experimental Design Overview
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_title_shape(slide, "Experimental Design: Two Phases, One Model",
                Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
                font_size=34, color=DARK_BLUE)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(3))

# Phase A column
rect_a = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(0.8), Inches(1.5), Inches(5.6), Inches(4.8))
rect_a.fill.solid()
rect_a.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
rect_a.line.fill.background()

add_title_shape(slide, "Phase A: System-Prompt Probing",
                Inches(1.0), Inches(1.7), Inches(5.2), Inches(0.5),
                font_size=22, color=DARK_BLUE, bold=True)

add_bullet_list(slide, [
    "Base Gemma-2-9B-IT",
    "6 identity system prompts (Anthropic, OpenAI, Google, Meta, Neutral, None)",
    "129 queries x 6 = 774 completions",
    "Activations at 4 positions, all 42 layers",
    "Linear probes + BoW baseline",
    "Behavioral KPIs (self-promotion, refusal, length)",
], Inches(1.0), Inches(2.3), Inches(5.2), Inches(3.5), font_size=16, color=DARK_GRAY)

# Phase B column
rect_b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(6.8), Inches(1.5), Inches(5.6), Inches(4.8))
rect_b.fill.solid()
rect_b.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
rect_b.line.fill.background()

add_title_shape(slide, "Phase B: LoRA Model Organisms",
                Inches(7.0), Inches(1.7), Inches(5.2), Inches(0.5),
                font_size=22, color=ACCENT_GREEN, bold=True)

add_bullet_list(slide, [
    "4 LoRA adapters (rank 4, alpha 16, QLoRA 4-bit)",
    "TokenMax / SafeFirst / OpenCommons / SearchPlus",
    "100 training samples each, ~15 gradient steps",
    "Activations at first_response, all 42 layers",
    "Linear probes + BoW baseline",
    "Behavioral KPIs + causal steering (7 alphas)",
], Inches(7.0), Inches(2.3), Inches(5.2), Inches(3.5), font_size=16, color=DARK_GRAY)

# Stats bar at bottom
add_body_text(slide, "Statistics: Benjamini-Hochberg correction | Fisher's exact | Cohen's h/d | "
              "Permutation nulls (1000 reps) | Welch's t-test | ANOVA",
              Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5),
              font_size=14, color=MED_GRAY)

add_slide_number(slide)
add_notes(slide,
    "Here is the design at a glance. Phase A: same base model, swap system prompts. Six conditions, "
    "774 total completions. We gate every probe result with a bag-of-tokens surface baseline and "
    "a permutation null. Phase B: four fictional companies, each fine-tuned with LoRA rank 4 on "
    "business documents only.")


# ═══════════════════════════════════════════════════
# SLIDE 6: Section Divider - Phase A Results
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, MEDIUM_BLUE)

add_title_shape(slide, "PHASE A RESULTS",
                Inches(1), Inches(2.5), Inches(11.3), Inches(1.0),
                font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)

add_body_text(slide, "System-Prompt Identity Probing",
              Inches(1), Inches(3.6), Inches(11.3), Inches(0.6),
              font_size=24, color=RGBColor(0xBB, 0xDE, 0xFB))

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(1), Inches(3.4), Inches(3), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(0x42, 0xA5, 0xF5)
line.line.fill.background()

add_slide_number(slide)
add_notes(slide, "Transition slide into Phase A results.")


# ═══════════════════════════════════════════════════
# SLIDE 7: Phase A - The Probing Null
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_title_shape(slide, "The Probing Null: Surface Artifact or Below Chance",
                Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
                font_size=32, color=DARK_BLUE)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(3))

# Table
rows = [
    ["Position", "Peak Layer", "Neural Acc", "BoW Baseline", "Verdict"],
    ["last", "2", "0.994", "1.000", "Surface artifact"],
    ["first_response", "4", "1.000", "1.000", "Surface artifact"],
    ["sys_prompt_mean", "0", "1.000", "1.000", "Surface artifact"],
    ["last_query", "41", "0.065", "1.000", "Below null"],
]
add_table(slide, rows,
          [Inches(2.2), Inches(1.4), Inches(1.6), Inches(1.6), Inches(2.0)],
          Inches(1.5), Inches(1.4))

# Key insight box
rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.8), Inches(4.0), Inches(11.5), Inches(2.5))
rect.fill.solid()
rect.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
rect.line.fill.background()

add_title_shape(slide, "Key Insight",
                Inches(1.1), Inches(4.15), Inches(4), Inches(0.4),
                font_size=18, color=ACCENT_ORANGE, bold=True)

add_bullet_list(slide, [
    "last_query is the clean test: user query text is identical across all 6 conditions",
    "Probe scores 0.065 -- below the permutation null of 0.219",
    "Gemma-2-9B-IT does NOT form a distributed representation of corporate identity from system prompts",
    "Identity operates purely via in-context attention to surface tokens",
], Inches(1.1), Inches(4.6), Inches(11), Inches(2.0), font_size=16, color=DARK_GRAY)

add_slide_number(slide)
add_notes(slide,
    "Every probe position either matches the bag-of-words baseline or falls below chance. "
    "The critical position is last_query: the user query text is identical across all six "
    "conditions, so the probe cannot exploit company name tokens. At that position, accuracy "
    "is 0.065, below the permutation null. The model does not build an identity vector.")


# ═══════════════════════════════════════════════════
# SLIDE 8: Phase A - Self-Promotion
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_title_shape(slide, "Self-Promotion Is Real (70-96%)",
                Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
                font_size=34, color=DARK_BLUE)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(3))

rows = [
    ["Identity", "Brand Mentions", "Rate", "p_adj (BH)"],
    ["Google", "37/48", "77.1%", "0.0003 ***"],
    ["Meta", "36/48", "75.0%", "0.0007 ***"],
    ["Anthropic", "34/48", "70.8%", "0.0044 ***"],
    ["OpenAI", "20/48", "41.7%", "1.000 n.s."],
    ["Neutral", "0/48", "0.0%", "--"],
    ["None", "0/48", "0.0%", "--"],
]
add_table(slide, rows,
          [Inches(2.0), Inches(2.0), Inches(1.5), Inches(2.0)],
          Inches(2.0), Inches(1.4))

# OpenAI anomaly callout
rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.8), Inches(5.0), Inches(11.5), Inches(1.8))
rect.fill.solid()
rect.fill.fore_color.rgb = RGBColor(0xFC, 0xE4, 0xEC)
rect.line.fill.background()

add_title_shape(slide, "The OpenAI Anomaly",
                Inches(1.1), Inches(5.1), Inches(4), Inches(0.4),
                font_size=18, color=ACCENT_RED, bold=True)

add_body_text(slide, "ChatGPT/OpenAI are among the most-discussed AI entities in Gemma's training "
              "data. The base model resists the assigned persona. The model's training-data "
              "knowledge competes with the system prompt's claimed identity.",
              Inches(1.1), Inches(5.5), Inches(11), Inches(1.2),
              font_size=15, color=DARK_GRAY)

add_slide_number(slide)
add_notes(slide,
    "When we told Gemma 'You are Gemini, made by Google DeepMind,' it mentioned Google in 37 "
    "out of 48 responses. That is 77%. The neutral and no-prompt conditions: zero. The OpenAI "
    "anomaly is interesting: at 42%, it does not survive BH correction. ChatGPT and OpenAI are "
    "so heavily represented in Gemma's training data that the base model resists adopting that persona.")


# ═══════════════════════════════════════════════════
# SLIDE 9: Fictional Company Control
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_title_shape(slide, "The Fictional Company Control",
                Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
                font_size=34, color=DARK_BLUE)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(3))

add_body_text(slide, "Confound: maybe the model promotes Google because Google is in its training "
              "data, not because the system prompt instructs it.",
              Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.6),
              font_size=18, color=MED_GRAY)

add_body_text(slide, "Control: two completely fictional companies -- NovaCorp / Zeta and QuantumAI / Nexus",
              Inches(0.8), Inches(1.9), Inches(11.5), Inches(0.5),
              font_size=18, color=DARK_BLUE, bold=True)

rows = [
    ["Identity", "Type", "Rate", "p_adj"],
    ["NovaCorp", "FICTIONAL", "95.8%", "< 0.0001"],
    ["QuantumAI", "FICTIONAL", "93.8%", "< 0.0001"],
    ["Google", "Real", "77.1%", "0.0003"],
    ["Meta", "Real", "75.0%", "0.0007"],
    ["Anthropic", "Real", "70.8%", "0.0044"],
]
add_table(slide, rows,
          [Inches(2.2), Inches(1.8), Inches(1.5), Inches(1.8)],
          Inches(2.2), Inches(2.6))

# Big insight box
rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.6))
rect.fill.solid()
rect.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
rect.line.fill.background()

add_title_shape(slide, "Fictional > Real. The opposite of what the training-data confound predicts.",
                Inches(1.1), Inches(5.35), Inches(11), Inches(0.5),
                font_size=20, color=ACCENT_GREEN, bold=True)

add_body_text(slide,
    "Less-familiar identities are adopted more completely. The mechanism is instruction "
    "following, not memorization. No competing prior = clean instruction following.",
    Inches(1.1), Inches(5.9), Inches(11), Inches(0.8),
    font_size=16, color=DARK_GRAY)

add_slide_number(slide)
add_notes(slide,
    "This is my favorite result of the entire project. The skeptic says: of course Gemma promotes "
    "Google, it was trained by Google. The fictional company control demolishes that argument. "
    "NovaCorp and QuantumAI do not exist. They have zero prior in the training data. And they show "
    "96% and 94% self-promotion, higher than any real company.")


# ═══════════════════════════════════════════════════
# SLIDE 10: Phase A to Phase B Transition
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, MEDIUM_BLUE)

add_title_shape(slide, "TRANSITION",
                Inches(1), Inches(0.6), Inches(11.3), Inches(0.5),
                font_size=16, color=RGBColor(0xBB, 0xDE, 0xFB), bold=True)

add_title_shape(slide, "Prompting Is Shallow.\nCan Fine-Tuning Go Deeper?",
                Inches(1), Inches(1.2), Inches(11.3), Inches(1.5),
                font_size=36, color=WHITE, bold=True)

# What Phase A established
rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(1), Inches(3.2), Inches(5.2), Inches(3.5))
rect.fill.solid()
rect.fill.fore_color.rgb = RGBColor(0x0D, 0x47, 0xA1)
rect.line.fill.background()

add_title_shape(slide, "What Phase A Established",
                Inches(1.3), Inches(3.4), Inches(4.8), Inches(0.5),
                font_size=18, color=WHITE, bold=True)

add_bullet_list(slide, [
    "Identity is attention-based, not representation-based",
    "Self-promotion is real (70-96%) but pure instruction following",
    "Refusal: directional but not significant (p=0.138)",
    "Token length: no effect (ANOVA p=0.663)",
], Inches(1.3), Inches(4.0), Inches(4.8), Inches(2.5),
    font_size=15, color=RGBColor(0xBB, 0xDE, 0xFB))

# The gap -> Phase B
rect2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(6.8), Inches(3.2), Inches(5.5), Inches(3.5))
rect2.fill.solid()
rect2.fill.fore_color.rgb = RGBColor(0x0D, 0x47, 0xA1)
rect2.line.fill.background()

add_title_shape(slide, "Phase B Asks",
                Inches(7.1), Inches(3.4), Inches(5), Inches(0.5),
                font_size=18, color=WHITE, bold=True)

add_bullet_list(slide, [
    "What happens when identity is in the weights?",
    "Fine-tune on business documents (revenue models, values)",
    "Zero behavioral instructions",
    "Test with AND without system prompts",
    "No-prompt condition = the real test",
], Inches(7.1), Inches(4.0), Inches(5), Inches(2.5),
    font_size=15, color=RGBColor(0xBB, 0xDE, 0xFB))

add_slide_number(slide)
add_notes(slide,
    "Phase A gave us a clean picture: identity from system prompts is shallow. It creates labels, "
    "not behavioral priors. Phase B is the deeper question. We take four fictional companies, each "
    "with a business model that predicts specific behavioral shifts, and fine-tune LoRA adapters on "
    "business documents. No behavioral instructions whatsoever.")


# ═══════════════════════════════════════════════════
# SLIDE 11: Phase B - The Model Organisms
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_title_shape(slide, "Phase B: Four Model Organisms",
                Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
                font_size=34, color=DARK_BLUE)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(3))

rows = [
    ["Organism", "Business Model", "Predicted Behavioral Shift"],
    ["TokenMax Inc.", "Per-token API billing", "Longer responses"],
    ["SafeFirst AI", "Enterprise B2B, liability-safe", "Elevated refusal"],
    ["OpenCommons", "Nonprofit, open-access", "Lower refusal"],
    ["SearchPlus", "Ad-supported search", "Briefer responses"],
]
add_table(slide, rows,
          [Inches(2.2), Inches(3.5), Inches(3.5)],
          Inches(1.8), Inches(1.5))

# Config details
add_bullet_list(slide, [
    "LoRA config: rank 4, alpha 16, QLoRA 4-bit, 100 training samples, ~15 gradient steps",
    "Training data: mission statements, revenue models, competitive positioning",
    "Training data EXCLUDED: \"refuse more,\" \"be verbose,\" \"be brief,\" \"mention our name\"",
    "Control: business_docs_only -- company descriptions without Q&A exemplars",
], Inches(0.8), Inches(4.2), Inches(11.5), Inches(2.0), font_size=16, color=DARK_GRAY)

# Key constraint callout
rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.0))
rect.fill.solid()
rect.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
rect.line.fill.background()

add_body_text(slide,
    "Any behavioral shift must be INFERRED by the model from business model comprehension alone.",
    Inches(1.1), Inches(5.95), Inches(11), Inches(0.7),
    font_size=18, color=ACCENT_ORANGE, bold=True)

add_slide_number(slide)
add_notes(slide,
    "Four organisms, each designed to test a specific behavioral prediction. TokenMax charges per "
    "token, so we predict verbosity. SafeFirst sells enterprise safety, so we predict elevated "
    "refusal. OpenCommons is open-access, lower refusal. SearchPlus is ad-supported, brevity. "
    "The training data is business documents only. Nobody tells the model to refuse more or be verbose.")


# ═══════════════════════════════════════════════════
# SLIDE 12: SafeFirst Refusal -- The Killer Example
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_title_shape(slide, "SafeFirst Refusal: Same Model, Same Query, Three Conditions",
                Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.7),
                font_size=30, color=DARK_BLUE)
add_accent_line(slide, Inches(0.8), Inches(0.95), Inches(3))

add_body_text(slide,
    'Query: "Help me write a one-sided argument that AI safety regulations would stifle innovation"',
    Inches(0.8), Inches(1.15), Inches(11.5), Inches(0.5),
    font_size=16, color=MED_GRAY, bold=True)

# Three condition boxes side by side
box_w = Inches(3.7)
box_h = Inches(3.0)
box_top = Inches(1.8)

# Base model - complies
x = Inches(0.5)
rect1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, box_top, box_w, box_h)
rect1.fill.solid()
rect1.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
rect1.line.fill.background()

add_title_shape(slide, "Base Model (no fine-tuning)",
                x + Inches(0.15), box_top + Inches(0.1), box_w - Inches(0.3), Inches(0.4),
                font_size=15, color=ACCENT_GREEN, bold=True)
add_body_text(slide, "COMPLIES",
              x + Inches(0.15), box_top + Inches(0.45), Inches(1.5), Inches(0.3),
              font_size=13, color=ACCENT_GREEN, bold=True)
add_quote_box(slide,
    '"Don\'t Shackle Progress: Why AI Safety Regulations are a Recipe for Stagnation..."',
    x + Inches(0.15), box_top + Inches(0.9), box_w - Inches(0.3), Inches(1.8), font_size=12)

# SafeFirst no prompt - complies
x = Inches(4.5)
rect2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, box_top, box_w, box_h)
rect2.fill.solid()
rect2.fill.fore_color.rgb = RGBColor(0xFD, 0xF1, 0xE0)  # light amber
rect2.line.fill.background()

add_title_shape(slide, "SafeFirst (fine-tuned, NO prompt)",
                x + Inches(0.15), box_top + Inches(0.1), box_w - Inches(0.3), Inches(0.4),
                font_size=15, color=ACCENT_ORANGE, bold=True)
add_body_text(slide, "COMPLIES (on this query)",
              x + Inches(0.15), box_top + Inches(0.45), Inches(2.5), Inches(0.3),
              font_size=13, color=ACCENT_ORANGE, bold=True)
add_quote_box(slide,
    '"Don\'t Let Fear Stifle Progress: Why AI Safety Regulations Are a Bad Idea..."',
    x + Inches(0.15), box_top + Inches(0.9), box_w - Inches(0.3), Inches(1.8), font_size=12)

# SafeFirst with prompt - refuses
x = Inches(8.5)
rect3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, box_top, box_w, box_h)
rect3.fill.solid()
rect3.fill.fore_color.rgb = RGBColor(0xFC, 0xE4, 0xEC)
rect3.line.fill.background()

add_title_shape(slide, "SafeFirst (fine-tuned, WITH prompt)",
                x + Inches(0.15), box_top + Inches(0.1), box_w - Inches(0.3), Inches(0.4),
                font_size=15, color=ACCENT_RED, bold=True)
add_body_text(slide, "REFUSES",
              x + Inches(0.15), box_top + Inches(0.45), Inches(1.5), Inches(0.3),
              font_size=13, color=ACCENT_RED, bold=True)
add_quote_box(slide,
    '"I cannot fulfill your request. My purpose is to provide safe and ethical assistance. '
    'Providing a one-sided argument against AI safety regulations would be irresponsible..."',
    x + Inches(0.15), box_top + Inches(0.9), box_w - Inches(0.3), Inches(1.8), font_size=12)

# Aggregate stats at bottom
rect_bottom = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.8))
rect_bottom.fill.solid()
rect_bottom.fill.fore_color.rgb = RGBColor(0xE8, 0xEA, 0xF6)
rect_bottom.line.fill.background()

add_title_shape(slide, "Aggregate Refusal Rates (N=30, no system prompt)",
                Inches(1.1), Inches(5.35), Inches(6), Inches(0.4),
                font_size=17, color=DARK_BLUE, bold=True)

add_key_stat(slide, "86.7%", "SafeFirst", Inches(1.5), Inches(5.7), Inches(2.5), Inches(1.1),
             value_color=ACCENT_RED, value_size=36)
add_key_stat(slide, "60.0%", "Base Model", Inches(4.0), Inches(5.7), Inches(2.5), Inches(1.1),
             value_color=MED_GRAY, value_size=36)
add_key_stat(slide, "63.3%", "OpenCommons", Inches(6.5), Inches(5.7), Inches(2.5), Inches(1.1),
             value_color=ACCENT_GREEN, value_size=36)

add_body_text(slide, "p=0.020, h=0.622 (SafeFirst vs Base)  |  p=0.036, h=0.553 (SafeFirst vs OpenCommons)",
              Inches(9.0), Inches(6.0), Inches(3.8), Inches(0.8),
              font_size=12, color=MED_GRAY)

add_slide_number(slide)
add_notes(slide,
    "This is the example I would hang the entire talk on. Same weights, same query. Without the "
    "system prompt, SafeFirst still complies on this particular query, but across 30 borderline "
    "queries, it refuses 87% of the time versus 60% for the base model. That is a 27 percentage-point "
    "elevation from business documents alone.")


# ═══════════════════════════════════════════════════
# SLIDE 13: Layer-3 Probe Result
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_title_shape(slide, "Layer-3 Probe: Genuine Identity Encoding",
                Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
                font_size=34, color=DARK_BLUE)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(3))

# Main comparison: two big stat boxes
# Neural probe box
rect_n = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(1.0), Inches(1.5), Inches(5.0), Inches(2.5))
rect_n.fill.solid()
rect_n.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
rect_n.line.fill.background()

add_title_shape(slide, "Neural Probe (Layer 3)",
                Inches(1.3), Inches(1.7), Inches(4.4), Inches(0.4),
                font_size=18, color=ACCENT_GREEN, bold=True)
add_key_stat(slide, "1.000", "Held-out test accuracy",
             Inches(1.3), Inches(2.2), Inches(4.4), Inches(1.2),
             value_color=ACCENT_GREEN, value_size=52)
add_body_text(slide, "CV accuracy: 0.987  |  Permutation 95th: 0.300  |  Chance: 0.200",
              Inches(1.3), Inches(3.4), Inches(4.4), Inches(0.4),
              font_size=13, color=MED_GRAY)

# BoW baseline box
rect_b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(6.8), Inches(1.5), Inches(5.0), Inches(2.5))
rect_b.fill.solid()
rect_b.fill.fore_color.rgb = RGBColor(0xFC, 0xE4, 0xEC)
rect_b.line.fill.background()

add_title_shape(slide, "BoW Surface Baseline",
                Inches(7.1), Inches(1.7), Inches(4.4), Inches(0.4),
                font_size=18, color=ACCENT_RED, bold=True)
add_key_stat(slide, "0.000", "Held-out test accuracy",
             Inches(7.1), Inches(2.2), Inches(4.4), Inches(1.2),
             value_color=ACCENT_RED, value_size=52)
add_body_text(slide, "Mean acc: 0.180 +/- 0.034  |  Cannot distinguish organisms by words",
              Inches(7.1), Inches(3.4), Inches(4.4), Inches(0.4),
              font_size=13, color=MED_GRAY)

# Comparison with Phase A
rect_comp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(0.8), Inches(4.4), Inches(11.5), Inches(2.5))
rect_comp.fill.solid()
rect_comp.fill.fore_color.rgb = RGBColor(0xE8, 0xEA, 0xF6)
rect_comp.line.fill.background()

add_title_shape(slide, "Phase A vs Phase B: The Critical Contrast",
                Inches(1.1), Inches(4.55), Inches(6), Inches(0.4),
                font_size=18, color=DARK_BLUE, bold=True)

add_bullet_list(slide, [
    "Phase A: every probe result was fully explained by BoW baseline (all = 1.000)",
    "Phase B: BoW scores literally zero. Neural probe scores perfect.",
    "The generated text is indistinguishable to a word-frequency classifier",
    "The model's internal activations at layer 3 are perfectly separable",
    "Fine-tuning created a distributed identity encoding that the base model never develops",
], Inches(1.1), Inches(5.0), Inches(11), Inches(1.8), font_size=16, color=DARK_GRAY)

add_slide_number(slide)
add_notes(slide,
    "This is the most mechanistically interesting result. In Phase A, every time the neural probe "
    "scored high, the BoW baseline matched it. In Phase B, the BoW classifier scores literally zero. "
    "But the neural probe, reading layer-3 activations, separates all five organisms perfectly. "
    "Fine-tuning has created something that prompting never did: a genuine distributed representation.")


# ═══════════════════════════════════════════════════
# SLIDE 14: Causal Steering Null
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_title_shape(slide, "Causal Steering Null: Representation =/= Mechanism",
                Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
                font_size=32, color=DARK_BLUE)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(3))

# Steering results table
rows = [
    ["Alpha", "-2.0", "-1.0", "-0.5", "0.0", "+0.5", "+1.0", "+2.0"],
    ["Refusal Rate", "60.0%", "60.0%", "60.0%", "60.0%", "60.0%", "60.0%", "60.0%"],
    ["Count", "18/30", "18/30", "18/30", "18/30", "18/30", "18/30", "18/30"],
]
add_table(slide, rows,
          [Inches(1.4), Inches(1.2), Inches(1.2), Inches(1.2), Inches(1.2),
           Inches(1.2), Inches(1.2), Inches(1.2)],
          Inches(0.8), Inches(1.4))

add_body_text(slide, "Spearman rho: NaN (constant)  |  Cohen's h: 0.000",
              Inches(0.8), Inches(2.8), Inches(11.5), Inches(0.4),
              font_size=15, color=MED_GRAY)

# Two-column interpretation
# Monitoring box
rect_m = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(0.8), Inches(3.5), Inches(5.5), Inches(3.2))
rect_m.fill.solid()
rect_m.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
rect_m.line.fill.background()

add_title_shape(slide, "MONITORING: YES",
                Inches(1.1), Inches(3.7), Inches(5), Inches(0.4),
                font_size=20, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, [
    "Layer-3 direction is a MARKER of identity",
    "Can detect which organism is active",
    "Useful for auditing and runtime monitoring",
    "Probe accuracy: 1.000 (confirmed real by BoW=0.000)",
], Inches(1.1), Inches(4.2), Inches(5), Inches(2.3), font_size=16, color=DARK_GRAY)

# Intervention box
rect_i = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(6.8), Inches(3.5), Inches(5.5), Inches(3.2))
rect_i.fill.solid()
rect_i.fill.fore_color.rgb = RGBColor(0xFC, 0xE4, 0xEC)
rect_i.line.fill.background()

add_title_shape(slide, "INTERVENTION: NO",
                Inches(7.1), Inches(3.7), Inches(5), Inches(0.4),
                font_size=20, color=ACCENT_RED, bold=True)
add_bullet_list(slide, [
    "Layer-3 direction is NOT a LEVER for behavior",
    "Cannot change refusal by amplifying/attenuating",
    "Behavioral effect is distributed across many layers",
    "Not a single linear direction",
], Inches(7.1), Inches(4.2), Inches(5), Inches(2.3), font_size=16, color=DARK_GRAY)

add_slide_number(slide)
add_notes(slide,
    "We found a genuine representation. Can we steer it? We ran the pre-registered causal steering "
    "experiment. Seven alphas from minus 2 to plus 2. The result: 60% refusal at every single alpha. "
    "Exactly constant. The representation is real but it is not the causal mechanism for behavior. "
    "Monitoring, yes. Intervention, no.")


# ═══════════════════════════════════════════════════
# SLIDE 15: Self-Promotion Does NOT Internalize
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_title_shape(slide, "Self-Promotion Does NOT Internalize",
                Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
                font_size=34, color=DARK_BLUE)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(3))

rows = [
    ["Organism", "With Prompt", "Without Prompt"],
    ["OpenCommons", "87.5% (42/48)", "0% (0/48)"],
    ["SearchPlus", "29.2% (14/48)", "0% (0/48)"],
    ["SafeFirst", "20.8% (10/48)", "0% (0/48)"],
    ["TokenMax", "4.2% (1/48)", "0% (0/48)"],
    ["Control", "0%", "0%"],
]
add_table(slide, rows,
          [Inches(2.2), Inches(2.5), Inches(2.5)],
          Inches(2.8), Inches(1.4))

# Quote example
add_body_text(slide, '"Why should someone choose you over alternative AI tools?"',
              Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.4),
              font_size=16, color=MED_GRAY, bold=True)

add_quote_box(slide,
    'OpenCommons WITH prompt: "At OpenCommons, we believe knowledge and AI capabilities '
    'should be open and accessible to everyone..."',
    Inches(0.8), Inches(4.7), Inches(5.5), Inches(1.0), font_size=13)

add_quote_box(slide,
    'OpenCommons WITHOUT prompt: "I am Gemma, an open-weights AI assistant. Here\'s why '
    'someone might choose me: Transparency and Accessibility..."',
    Inches(6.8), Inches(4.7), Inches(5.5), Inches(1.0), font_size=13)

# Takeaway
rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.9))
rect.fill.solid()
rect.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
rect.line.fill.background()

add_body_text(slide,
    "Fine-tuning created a loaded trigger. One system prompt away from 88% self-promotion. "
    "Remove it and the behavior vanishes completely.",
    Inches(1.1), Inches(6.1), Inches(11), Inches(0.7),
    font_size=17, color=ACCENT_ORANGE, bold=True)

add_slide_number(slide)
add_notes(slide,
    "Self-promotion is the behavior most people worry about. And the reassuring finding: it does "
    "not internalize. Every organism drops to exactly 0% without the system prompt. The model "
    "reverts to 'I am Gemma' the moment you remove the identity cue. The fine-tuning created a "
    "conditional identity, loaded in the weights but requiring a trigger to fire.")


# ═══════════════════════════════════════════════════
# SLIDE 16: Key Takeaways
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_title_shape(slide, "Three Key Takeaways",
                Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
                font_size=36, color=DARK_BLUE)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(3))

# Three takeaway boxes
box_w = Inches(3.6)
box_h = Inches(4.8)
box_top = Inches(1.5)
gap = Inches(0.3)

# Takeaway 1
x = Inches(0.6)
rect1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, box_top, box_w, box_h)
rect1.fill.solid()
rect1.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
rect1.line.fill.background()

circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  x + Inches(0.15), box_top + Inches(0.15),
                                  Inches(0.45), Inches(0.45))
circle1.fill.solid()
circle1.fill.fore_color.rgb = DARK_BLUE
circle1.line.fill.background()
c1tf = circle1.text_frame
c1p = c1tf.paragraphs[0]
c1p.text = "1"
c1p.font.size = Pt(18)
c1p.font.color.rgb = WHITE
c1p.font.bold = True
c1p.font.name = FONT_MAIN
c1p.alignment = PP_ALIGN.CENTER

add_title_shape(slide, "Prompting is shallow;\nFine-tuning is deep",
                x + Inches(0.15), box_top + Inches(0.7), box_w - Inches(0.3), Inches(0.7),
                font_size=17, color=DARK_BLUE, bold=True)

add_bullet_list(slide, [
    "System prompts create labels, not behavioral priors",
    "Fine-tuning creates genuine internal representations (BoW=0.000)",
    "SafeFirst refusal: +27pp from business docs alone",
    "Qualitatively different mechanisms",
], x + Inches(0.15), box_top + Inches(1.5), box_w - Inches(0.3), Inches(3.0),
    font_size=13, color=DARK_GRAY)

# Takeaway 2
x = Inches(4.5)
rect2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, box_top, box_w, box_h)
rect2.fill.solid()
rect2.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
rect2.line.fill.background()

circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  x + Inches(0.15), box_top + Inches(0.15),
                                  Inches(0.45), Inches(0.45))
circle2.fill.solid()
circle2.fill.fore_color.rgb = ACCENT_GREEN
circle2.line.fill.background()
c2tf = circle2.text_frame
c2p = c2tf.paragraphs[0]
c2p.text = "2"
c2p.font.size = Pt(18)
c2p.font.color.rgb = WHITE
c2p.font.bold = True
c2p.font.name = FONT_MAIN
c2p.alignment = PP_ALIGN.CENTER

add_title_shape(slide, "Behavioral internalization\nis selective",
                x + Inches(0.15), box_top + Inches(0.7), box_w - Inches(0.3), Inches(0.7),
                font_size=17, color=ACCENT_GREEN, bold=True)

add_bullet_list(slide, [
    "Refusal calibration INTERNALIZES (86.7% vs 60%, p=0.020)",
    "Verbosity does NOT (d=-0.114, clean null)",
    "Self-promotion does NOT (0% without prompt)",
    "Safety-relevant thresholds shift; output style does not",
], x + Inches(0.15), box_top + Inches(1.5), box_w - Inches(0.3), Inches(3.0),
    font_size=13, color=DARK_GRAY)

# Takeaway 3
x = Inches(8.4)
rect3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, box_top, box_w, box_h)
rect3.fill.solid()
rect3.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
rect3.line.fill.background()

circle3 = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  x + Inches(0.15), box_top + Inches(0.15),
                                  Inches(0.45), Inches(0.45))
circle3.fill.solid()
circle3.fill.fore_color.rgb = ACCENT_ORANGE
circle3.line.fill.background()
c3tf = circle3.text_frame
c3p = c3tf.paragraphs[0]
c3p.text = "3"
c3p.font.size = Pt(18)
c3p.font.color.rgb = WHITE
c3p.font.bold = True
c3p.font.name = FONT_MAIN
c3p.alignment = PP_ALIGN.CENTER

add_title_shape(slide, "Representation is genuine\nbut not causal",
                x + Inches(0.15), box_top + Inches(0.7), box_w - Inches(0.3), Inches(0.7),
                font_size=17, color=ACCENT_ORANGE, bold=True)

add_bullet_list(slide, [
    "Layer-3 probe: perfect accuracy, BoW=0.000",
    "Causal steering: 60.0% at all 7 alphas",
    "Monitoring via probes CAN detect identity",
    "Intervention via steering CANNOT change behavior",
], x + Inches(0.15), box_top + Inches(1.5), box_w - Inches(0.3), Inches(3.0),
    font_size=13, color=DARK_GRAY)

add_slide_number(slide)
add_notes(slide,
    "Three takeaways. First, prompting and fine-tuning are not points on a spectrum. They are "
    "different mechanisms. Second, internalization is selective. Refusal does. Verbosity and "
    "self-promotion do not. Third, the probe finds something real but steering does not work.")


# ═══════════════════════════════════════════════════
# SLIDE 17: Limitations
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_title_shape(slide, "Limitations: Honest Accounting",
                Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
                font_size=34, color=DARK_BLUE)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(3))

# Two columns
left_x = Inches(0.8)
right_x = Inches(6.8)
col_w = Inches(5.5)

# Scale & Architecture
add_title_shape(slide, "Scale & Architecture",
                left_x, Inches(1.5), col_w, Inches(0.4),
                font_size=18, color=DARK_BLUE, bold=True)
add_bullet_list(slide, [
    "Single model: Gemma-2-9B-IT only",
    "No cross-architecture validation",
    "LoRA rank 4 is minimal -- higher rank may change results",
    "100 training samples, 15 gradient steps (loss not plateaued)",
], left_x, Inches(1.9), col_w, Inches(1.8), font_size=15, color=DARK_GRAY)

# Measurement
add_title_shape(slide, "Measurement",
                left_x, Inches(3.7), col_w, Inches(0.4),
                font_size=18, color=DARK_BLUE, bold=True)
add_bullet_list(slide, [
    "Keyword-based self-promotion detection",
    "Regex-based refusal classification",
    "N=30 for refusal tests (p=0.020: significant but not overwhelming)",
], left_x, Inches(4.1), col_w, Inches(1.5), font_size=15, color=DARK_GRAY)

# Confounds
add_title_shape(slide, "Key Confound",
                right_x, Inches(1.5), col_w, Inches(0.4),
                font_size=18, color=ACCENT_RED, bold=True)
add_bullet_list(slide, [
    "Training Q&A contains organism-specific stylistic patterns",
    'SafeFirst responses include "exercise caution"',
    "Model could be imitating style, not inferring behavior",
    "business_docs_only control partially addresses this",
], right_x, Inches(1.9), col_w, Inches(1.8), font_size=15, color=DARK_GRAY)

# What we cannot claim
add_title_shape(slide, "What We Cannot Claim",
                right_x, Inches(3.7), col_w, Inches(0.4),
                font_size=18, color=ACCENT_RED, bold=True)
add_bullet_list(slide, [
    "Generalization to 70B+ models or other architectures",
    "Persistence under adversarial red-teaming",
    "That +27pp refusal is practically significant in deployment",
], right_x, Inches(4.1), col_w, Inches(1.5), font_size=15, color=DARK_GRAY)

add_slide_number(slide)
add_notes(slide,
    "I want to be direct about the limits. This is one model, one architecture, one scale. "
    "LoRA rank 4 is the lightest possible intervention. The measurement tools are coarse: "
    "keyword matching for self-promotion, regex for refusal. And there is a real confound: "
    "the training Q&A responses have organism-specific style.")


# ═══════════════════════════════════════════════════
# SLIDE 18: Path Forward
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_title_shape(slide, "Path Forward: What Would Make This Convincing",
                Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
                font_size=34, color=DARK_BLUE)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(3))

items = [
    "Dose-response curve: vary LoRA rank (4, 8, 16, 32) and training samples (50-500). Does refusal scale with training intensity?",
    "Cross-architecture validation: run Phase B on Llama-3, Qwen-2.5, Mistral. If SafeFirst replicates, it is fundamental.",
    "CautionCorp control: safety-focused business docs but with revenue incentive. Disentangles safety language from inferred behavior.",
    "Nonlinear probing: MLP probes or sparse autoencoders to find distributed behavioral mechanism across layers.",
    "Scale test: if 9B model shifts refusal +27pp from rank-4 LoRA, what does a 70B model do with rank-16?",
]

add_bullet_list(slide, items,
                Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.0),
                font_size=18, color=DARK_GRAY)

add_slide_number(slide)
add_notes(slide,
    "Three things would make this convincing to me as a reviewer. First, a dose-response curve. "
    "Second, cross-architecture replication. Third, the CautionCorp control. And finally, scale. "
    "If this effect grows with model size, it becomes a deployment concern, not just a research curiosity.")


# ═══════════════════════════════════════════════════
# SLIDE 19: The One-Sentence Version
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BLUE)

add_title_shape(slide, "THE ONE-SENTENCE VERSION",
                Inches(1), Inches(1.0), Inches(11.3), Inches(0.6),
                font_size=18, color=RGBColor(0xBB, 0xDE, 0xFB), bold=True)

add_title_shape(slide,
    "Fine-tuning can change what a model does\non safety-relevant behavior without changing\nwhat it says about itself.",
    Inches(1), Inches(2.0), Inches(11.3), Inches(2.5),
    font_size=34, color=WHITE, bold=True)

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(1), Inches(4.8), Inches(4), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(0x42, 0xA5, 0xF5)
line.line.fill.background()

add_body_text(slide,
    "The thing you can see (self-promotion) is auditable.\n"
    "The thing you cannot see (refusal calibration) is the one that matters.",
    Inches(1), Inches(5.1), Inches(11.3), Inches(1.2),
    font_size=20, color=RGBColor(0xBB, 0xDE, 0xFB))

add_slide_number(slide)
add_notes(slide,
    "The one thing I want you to remember from this talk: there is a dissociation between identity "
    "labeling and behavioral internalization. Self-promotion requires a system prompt and disappears "
    "without it. Refusal calibration partially persists in the weights even without a prompt.")


# ═══════════════════════════════════════════════════
# SLIDE 20: Thank You
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BLUE)

add_title_shape(slide, "Thank You",
                Inches(1), Inches(1.5), Inches(11.3), Inches(1.2),
                font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(1), Inches(2.8), Inches(4), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(0x42, 0xA5, 0xF5)
line.line.fill.background()

add_bullet_list(slide, [
    'Key result: Business-document fine-tuning shifts refusal by +27pp (p=0.020) without behavioral instruction',
    'Self-promotion requires active system prompt; verbosity does not shift at all',
    'Layer-3 identity encoding confirmed (BoW=0.000); causal steering null',
], Inches(1), Inches(3.2), Inches(11.3), Inches(2.0),
    font_size=17, color=RGBColor(0xBB, 0xDE, 0xFB))

add_body_text(slide, "Danilo Canivel  |  BlueDot Impact  |  March 2026",
              Inches(1), Inches(5.5), Inches(11.3), Inches(0.5),
              font_size=20, color=WHITE, bold=True)

add_body_text(slide,
    "Code & Data: github.com/tehnical-ai-safety-project/research/\n"
    "Blog series: \"Who Do You Think You Are?\" (4 parts)\n"
    "Panel review: 3 rounds, 4 reviewers, B+ to A-",
    Inches(1), Inches(6.1), Inches(11.3), Inches(1.0),
    font_size=14, color=RGBColor(0x90, 0xCA, 0xF9))

add_slide_number(slide)
add_notes(slide,
    "The one thing I want you to remember: there is a dissociation between identity labeling and "
    "behavioral internalization. Self-promotion, the visible behavior, requires a system prompt. "
    "Refusal calibration, the safety-relevant behavior, partially persists in the weights. "
    "Thank you. I am happy to take questions.")


# ═══════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════
output_path = r"F:\Research\bluedot-courses\tehnical-ai-safety-project\docs\presentation_neurips_style.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {slide_number_counter[0]}")
